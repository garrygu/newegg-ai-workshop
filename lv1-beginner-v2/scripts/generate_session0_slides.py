"""
Generate Session_0_Setup_and_AI_Fundamentals.pptx
for middle school students — simple, fun, visual.

Requires: pip install python-pptx
Run from: lv1-beginner-v2/
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Color palette ──────────────────────────────────────────────────
C_BG        = RGBColor(0x0F, 0x0F, 0x1A)   # Dark navy
C_ACCENT    = RGBColor(0x6C, 0x63, 0xFF)   # Purple
C_ACCENT2   = RGBColor(0x00, 0xD4, 0xFF)   # Cyan
C_GREEN     = RGBColor(0x00, 0xE6, 0x76)   # Green
C_YELLOW    = RGBColor(0xFF, 0xD6, 0x00)   # Yellow
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xCC, 0xFF)   # Light purple text

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_title_slide(prs, title, subtitle, emoji="🤖"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide, C_BG)

    # Big emoji
    add_textbox(slide, emoji,
                Inches(5.5), Inches(0.5), Inches(2), Inches(1.5),
                font_size=72, align=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, title,
                Inches(1), Inches(2.0), Inches(11), Inches(1.8),
                font_size=48, bold=True, color=C_ACCENT2,
                align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, subtitle,
                Inches(1.5), Inches(3.8), Inches(10), Inches(1.2),
                font_size=24, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, "🎓 Newegg YouthAI Workshop  •  Session 0",
                Inches(0), Inches(6.9), Inches(13.33), Inches(0.5),
                font_size=14, color=RGBColor(0x55,0x55,0x88),
                align=PP_ALIGN.CENTER)

def add_section_slide(prs, number, title, emoji):
    """Dark intro slide for each Part."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, RGBColor(0x0A, 0x0A, 0x28))

    add_textbox(slide, emoji,
                Inches(5.5), Inches(1.5), Inches(2), Inches(1.5),
                font_size=60, align=PP_ALIGN.CENTER)

    add_textbox(slide, number,
                Inches(1), Inches(2.8), Inches(11), Inches(0.7),
                font_size=18, color=C_ACCENT, align=PP_ALIGN.CENTER, bold=True)

    add_textbox(slide, title,
                Inches(1), Inches(3.4), Inches(11), Inches(1.2),
                font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def add_content_slide(prs, title, bullets, emoji="", accent=C_ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Title row
    if emoji:
        add_textbox(slide, emoji, Inches(0.3), Inches(0.15), Inches(1), Inches(0.8), font_size=32)

    add_textbox(slide, title,
                Inches(1.2), Inches(0.15), Inches(11.5), Inches(0.8),
                font_size=30, bold=True, color=accent)

    # Bullets
    top = Inches(1.1)
    for i, (icon, text) in enumerate(bullets):
        add_textbox(slide, f"{icon}  {text}",
                    Inches(0.5), top + i * Inches(0.88),
                    Inches(12.3), Inches(0.82),
                    font_size=22, color=C_WHITE)

def add_two_col_slide(prs, title, left_items, right_items,
                      left_header="", right_header="",
                      left_color=C_ACCENT, right_color=C_GREEN):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT2
    bar.line.fill.background()

    add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
                font_size=30, bold=True, color=C_ACCENT2)

    # Divider
    div = slide.shapes.add_shape(1, Inches(6.5), Inches(1.0), Inches(0.05), Inches(5.8))
    div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0x33,0x33,0x55)
    div.line.fill.background()

    for col, (header, items, color) in enumerate([
        (left_header, left_items, left_color),
        (right_header, right_items, right_color)
    ]):
        left_x = Inches(0.4) if col == 0 else Inches(6.8)
        if header:
            add_textbox(slide, header, left_x, Inches(1.0), Inches(5.8), Inches(0.6),
                        font_size=20, bold=True, color=color)
        for i, (icon, text) in enumerate(items):
            add_textbox(slide, f"{icon}  {text}",
                        left_x, Inches(1.6) + i * Inches(0.85),
                        Inches(5.8), Inches(0.8),
                        font_size=20, color=C_WHITE)

def add_code_slide(prs, title, code_lines, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_YELLOW
    bar.line.fill.background()

    add_textbox(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
                font_size=28, bold=True, color=C_YELLOW)

    # Code box bg
    code_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(1.0), Inches(12.5), Inches(4.8))
    code_bg.fill.solid(); code_bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    code_bg.line.color.rgb = RGBColor(0x33, 0x33, 0x55)

    code_text = "\n".join(code_lines)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(16)
    run.font.color.rgb = C_GREEN
    run.font.name = "Courier New"

    if note:
        add_textbox(slide, f"💡 {note}", Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.6),
                    font_size=18, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════
#  BUILD THE PRESENTATION
# ══════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── 1. TITLE SLIDE ────────────────────────────────────────────────
add_title_slide(prs,
    "Session 0: Setup & AI Fundamentals",
    "Welcome to the Newegg AI Workshop! 🚀\n"
    "Today: Python basics • AI types • Key concepts • Ethics",
    "🧭")

# ── 2. AGENDA ─────────────────────────────────────────────────────
add_content_slide(prs, "What We'll Cover Today", [
    ("📚", "Part 1 — Getting Started with Jupyter Notebooks"),
    ("🐍", "Part 2 — Python Basics (variables, loops, functions)"),
    ("🌐", "Part 3 — The 2025 AI Landscape (3 types of AI)"),
    ("🧠", "Part 4 — Core AI Concepts (training vs inference)"),
    ("⚖️", "Part 5 — AI Ethics & Responsible Use"),
    ("🎯", "Part 6 — Quiz + Challenge!"),
], emoji="📋", accent=C_ACCENT2)

# ── PART 1: JUPYTER ───────────────────────────────────────────────
add_section_slide(prs, "PART 1", "Getting Started with Jupyter", "📓")

add_content_slide(prs, "What is a Jupyter Notebook?", [
    ("📄", "Interactive document where you write AND run code"),
    ("🧱", "Made of CELLS — each cell is a block of code or text"),
    ("▶️",  "Run a cell: press  Shift + Enter"),
    ("💡", "See results instantly — great for experimenting!"),
    ("🌐", "Works in your browser — no special software needed"),
], emoji="📓", accent=C_ACCENT)

add_code_slide(prs, "Your First Line of Python! 🎉",
    ['print("🎉 Hello AI World! Welcome to the future!")'],
    note="Try it! Type this in a cell and press Shift + Enter")

# ── PART 2: PYTHON BASICS ─────────────────────────────────────────
add_section_slide(prs, "PART 2", "Python Basics for AI", "🐍")

add_two_col_slide(prs, "Variables — Storing Information",
    left_header="📦 What is a Variable?",
    left_items=[
        ("🏷️", "A named container that holds data"),
        ("🔢", "name = \"AI Student\"   (text)"),
        ("🔢", "age = 16               (number)"),
        ("✅", "is_excited = True      (yes/no)"),
    ],
    right_header="🗂️ Why Do We Need Them?",
    right_items=[
        ("🤖", "Store things AI needs to remember"),
        ("🔄", "Reuse values without retyping"),
        ("📊", "Feed data into AI models"),
        ("💡", "Think: labeled sticky notes 🗒️"),
    ],
    left_color=C_ACCENT, right_color=C_GREEN)

add_two_col_slide(prs, "Lists & Loops — Doing Things Repeatedly",
    left_header="📋 Lists",
    left_items=[
        ("📦", "Hold multiple items at once"),
        ("🛠️", 'ai_tools = ["ChatGPT", "DALL-E"]'),
        ("🔢", "Access by position: ai_tools[0]"),
        ("📏", "len(ai_tools) → how many items"),
    ],
    right_header="🔄 Loops",
    right_items=[
        ("🔁", "Repeat an action automatically"),
        ("🏋️", "AI trains by looping over data"),
        ("💻", 'for tool in ai_tools: print(tool)'),
        ("💡", "Think: a machine assembly line 🏭"),
    ],
    left_color=C_ACCENT2, right_color=C_YELLOW)

add_code_slide(prs, "Functions — Reusable Blocks of Code 🔧",
    [
        'def greet_student(name, favorite_ai):',
        '    return f"👋 Welcome, {name}! You love {favorite_ai}!"',
        '',
        '# Call the function',
        'print(greet_student("Alex", "ChatGPT"))',
        '',
        '# OUTPUT:  👋 Welcome, Alex! You love ChatGPT!',
    ],
    note="Functions = recipes 📖 — write once, use many times!")

add_code_slide(prs, "Conditionals — Making Decisions (like AI does!) 🤔",
    [
        'confidence = 0.85   # AI\'s certainty level (0 to 1)',
        '',
        'if confidence > 0.9:',
        '    print("🎯 AI is very confident!")',
        'elif confidence > 0.7:',
        '    print("👍 AI is fairly confident.")',
        'else:',
        '    print("❓ AI is guessing.")',
        '',
        '# OUTPUT:  👍 AI is fairly confident.',
    ],
    note="AI uses math like this millions of times per second!")

# ── PART 3: AI LANDSCAPE ──────────────────────────────────────────
add_section_slide(prs, "PART 3", "The 2025 AI Landscape", "🌐")

add_content_slide(prs, "Three Types of Modern AI", [
    ("🎨", "GENERATIVE — Creates NEW content   →  ChatGPT, DALL-E, Stable Diffusion"),
    ("🔮", "PREDICTIVE — Makes predictions from data  →  Weather AI, spam filters"),
    ("🤖", "AGENTIC    — Takes ACTIONS on its own  →  Self-driving cars, coding agents"),
], emoji="🌐", accent=C_ACCENT2)

add_two_col_slide(prs, "Generative AI — The Creator 🎨",
    left_header="✅ What it does",
    left_items=[
        ("📝", "Writes stories, code, essays"),
        ("🖼️", "Generates images from text"),
        ("🎵", "Composes music"),
        ("🎬", "Creates videos"),
    ],
    right_header="🛠️ Tools you'll use",
    right_items=[
        ("🤖", "ChatGPT / Gemini (Session 1)"),
        ("🎨", "Stable Diffusion (Session 2)"),
        ("💬", "Hugging Face (Session 4)"),
        ("💡", '"Create something from nothing"'),
    ],
    left_color=C_ACCENT, right_color=C_ACCENT2)

add_two_col_slide(prs, "Predictive & Agentic AI",
    left_header="🔮 Predictive AI",
    left_items=[
        ("☀️", "Predicts weather"),
        ("📧", "Detects spam emails"),
        ("🏷️", "Classifies images (Session 3)"),
        ("💡", '"What will happen next?"'),
    ],
    right_header="🤖 Agentic AI",
    right_items=[
        ("🚗", "Self-driving cars"),
        ("💻", "Coding agents (like Cursor)"),
        ("🎮", "Our AI Game! (Session 5)"),
        ("💡", '"Do something on my behalf"'),
    ],
    left_color=C_GREEN, right_color=C_YELLOW)

add_content_slide(prs, "AI Tools We'll Build With", [
    ("🧠", "Session 1 — ChatGPT / Gemini API  →  Build your own AI assistant"),
    ("🎨", "Session 2 — Stable Diffusion  →  Generate AI artwork"),
    ("🔍", "Session 3 — PyTorch + CIFAR-10  →  Train an image classifier"),
    ("💬", "Session 4 — Hugging Face + Whisper  →  Sentiment chatbot + voice"),
    ("🎮", "Session 5 — All together  →  Build an AI guessing game!"),
], emoji="🛠️", accent=C_YELLOW)

# ── PART 4: CORE AI CONCEPTS ──────────────────────────────────────
add_section_slide(prs, "PART 4", "Core AI Concepts", "🧠")

add_content_slide(prs, "The AI Learning Pipeline", [
    ("📊", "DATA     →  Collection of examples  (1,000 cat & dog photos)"),
    ("🏋️",  "TRAINING →  AI studies the data and finds patterns"),
    ("🧠", "MODEL    →  The 'brain' that stores what was learned"),
    ("🔮", "INFERENCE →  Use the model to make predictions on new data"),
], emoji="🔄", accent=C_ACCENT2)

add_two_col_slide(prs, "Training vs Inference — Study vs Test",
    left_header="📚 Training  (Studying)",
    left_items=[
        ("🖼️", "Feed AI thousands of images"),
        ("🔁", "AI looks for patterns in data"),
        ("⏱️", "Takes hours or days to complete"),
        ("💻", "Needs a powerful GPU"),
    ],
    right_header="📝 Inference  (Test Day!)",
    right_items=[
        ("🆕", "Give AI a NEW image it's never seen"),
        ("⚡", "AI answers in milliseconds"),
        ("☁️", "Can run in the cloud cheaply"),
        ("✅", "This is what we do in class!"),
    ],
    left_color=C_ACCENT, right_color=C_GREEN)

# ── PART 5: ETHICS ────────────────────────────────────────────────
add_section_slide(prs, "PART 5", "AI Ethics & Responsible Use", "⚖️")

add_content_slide(prs, "Things That Can Go Wrong with AI ⚠️", [
    ("😤", "BIAS — AI reflects unfair human patterns in its training data"),
    ("🤥", "HALLUCINATIONS — AI confidently makes up false information"),
    ("🔒", "PRIVACY — AI can be trained on your data without permission"),
    ("🎭", "DEEPFAKES — Realistic fake videos/images of real people"),
    ("💼", "JOB IMPACT — Some jobs may be affected by automation"),
], emoji="⚠️", accent=RGBColor(0xFF, 0x66, 0x33))

add_content_slide(prs, "How to Use AI Responsibly ✅", [
    ("🔍", "Always VERIFY — don't trust AI output blindly"),
    ("🚫", "Never use AI to harm, cheat, or deceive others"),
    ("📜", "Respect copyright — don't steal content with AI"),
    ("🤝", "Be TRANSPARENT — tell people when you used AI"),
    ("💭", "Think about the impact on others before you post"),
], emoji="✅", accent=C_GREEN)

# ── QUIZ ──────────────────────────────────────────────────────────
add_section_slide(prs, "QUIZ TIME!", "Test Your Knowledge", "🎯")

add_two_col_slide(prs, "Quick Quiz — Can You Answer These?",
    left_header="❓ Questions",
    left_items=[
        ("1️⃣", "What AI type CREATES new content?"),
        ("2️⃣", "Taking the test = Training or Inference?"),
        ("3️⃣", "When AI makes up false info it's called...?"),
        ("4️⃣", "What key runs a Jupyter cell?"),
    ],
    right_header="✅ Answers",
    right_items=[
        ("🎨", "Generative AI"),
        ("📝", "Inference"),
        ("🤥", "Hallucination"),
        ("⌨️", "Shift + Enter"),
    ],
    left_color=C_ACCENT2, right_color=C_GREEN)

# ── WRAP UP ───────────────────────────────────────────────────────
add_content_slide(prs, "What You Learned Today 🏆", [
    ("✅", "Jupyter Notebooks — how to write and run code interactively"),
    ("✅", "Python basics — variables, loops, functions, if/else"),
    ("✅", "3 types of AI — Generative, Predictive, Agentic"),
    ("✅", "AI pipeline — Data → Training → Model → Inference"),
    ("✅", "AI ethics — bias, hallucinations, responsible use"),
], emoji="🏆", accent=C_YELLOW)

add_content_slide(prs, "Coming Up Next — Session 1 🚀", [
    ("🧠", "How ChatGPT and Large Language Models actually work"),
    ("✍️",  "Prompt engineering — zero-shot, few-shot, chain-of-thought"),
    ("🐍", "Call the OpenAI API from Python code"),
    ("🤖", "Build your very own AI assistant!"),
    ("📖", "Pre-read: 'How ChatGPT Works' on YouTube (Computerphile)"),
], emoji="🔮", accent=C_ACCENT2)

# ── SAVE ──────────────────────────────────────────────────────────
out_path = "Session_0_Slides.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
